from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

# This function gonna set a gradient background with 2 color (one for the top of the page and one for the bottom)
def SetBackground(prs, slide, data):
    # We create a an image on all the background of the page
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        0,
        0,
        prs.slide_width,
        prs.slide_height
    )

    # Initialize the shape's fill with a gradient
    fill = shape.fill
    fill.gradient()  

    # Access the gradient stops to set start and end colors
    gradient_stops = fill.gradient_stops

    # We convert the hexadecimal color in rgb color
    color = tuple(int(data["Start"][i:i+2], 16) for i in (0, 2, 4))
    
    # We set the color to the top of the gradient
    stop_1 = gradient_stops[0]
    stop_1.position = 0.0
    stop_1.color.rgb = RGBColor(color[0], color[1], color[2])

    # We convert the hexadecimal color in rgb color
    color = tuple(int(data["End"][i:i+2], 16) for i in (0, 2, 4))

    # We set the color to the bottom of the gradient
    stop_2 = gradient_stops[1]
    stop_2.position = 1.0
    stop_2.color.rgb = RGBColor(color[0], color[1], color[2])


from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.util import Pt, Cm

# This function gonna add a text to a powerpoint page
def CreateText(data, slide):
    # We add a box to the text (and the text inside this box will break the line if it is too big)
    title = slide.shapes.add_textbox(Cm(data["Left"]), Cm(data["Top"]), Cm(data["Width"]), Cm(data["Height"]))
    text_frame = title.text_frame
    text_frame.word_wrap = True

    # We add a paragraph who gonna handle the align of the text
    paragraph = text_frame.add_paragraph()
    if (data["Align"] == "center"):
        paragraph.alignment = PP_ALIGN.CENTER
    if (data["Align"] == "left"):
        paragraph.alignment = PP_ALIGN.LEFT
    if (data["Align"] == "right"):
        paragraph.alignment = PP_ALIGN.RIGHT
    if (data["Align"] == "justify"):
        paragraph.alignment = PP_ALIGN.JUSTIFY

    # We convert the hexadecimal color to a rgb color
    color = tuple(int(data["FontColor"][i:i+2], 16) for i in (0, 2, 4))

    # We add an exectuable code (who gonna contain the text)
    run = paragraph.add_run()
    # We set the font parameters and the text
    run.font.name = data["FontName"]
    run.font.size = Pt(data["FontSize"])
    run.font.color.rgb = RGBColor(color[0], color[1], color[2])
    run.text = data["Text"]

    # We set the font in bold, underline, italic if necessary
    if (data["Bold"]):
        run.font.bold = True
    if (data["Italic"]):
        run.font.italic = True
    if (data["Underline"]):
        run.font.underline = True

# This function gonna add an image to a powerpoint page
def CreateImage(data, slide):
    # We try to set the width and height in cm (but since they can't be zero, the conversion may fail)
    try:
        width = Cm(data["Width"])
    except (ValueError, TypeError):
        width = None

    try:
        height = Cm(data["Height"])
    except (ValueError, TypeError):
        height = None

    # We add the picture to the page with all his data convert in Cm
    slide.shapes.add_picture(data["Path"], Cm(data["Left"]), Cm(data["Top"]), width, height)