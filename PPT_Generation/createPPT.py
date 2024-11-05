import json
from pptx import Presentation

from PPT_Generation.createElem import CreateText, SetBackground, CreateImage

# This function gonna create the powerpoint file
def CreatePPT(JsonPath):
    # We open the json file and store his data in a dict
    with open(JsonPath, 'r', encoding='utf-8') as f:
        data = json.load(f)

    # We initialize the powerpoint as a Presentation instance
    prs = Presentation()
    # We set the slide layout type as the 6 (a free slide)
    lyt = prs.slide_layouts[6]

    # We loop on the pages on the datas of the json
    for i in data:
        # We add a new slide
        slide = prs.slides.add_slide(lyt)

        # We loop on the element of the page
        for y in data[i]:
            # We set the background
            if y == "Background":
                SetBackground(prs, slide, data[i][y])
            # We set the image
            if "Path" in data[i][y]:
                CreateImage(data[i][y], slide)
            # We set the text
            if "Text" in data[i][y]:
                CreateText(data[i][y], slide)

    # We save the powerpoint
    prs.save("powerpoint.ppt")